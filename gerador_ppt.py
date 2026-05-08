"""
gerador_ppt.py  —  PPT idêntico ao template Sesc/IFec RJ
=========================================================
Estratégia: clona os slides do arquivo original como templates
e substitui apenas título, gráfico e nota de base.
Todos os grupos decorativos (freeforms, chevrons) são preservados
pixel a pixel do original.

Requer: Carnaval_Sapucaí_-_Sesc.pptx na mesma pasta (template base).
Se não encontrar, usa fallback com formas simples.
"""

from __future__ import annotations
import copy, os, re
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ─────────────────────────────────────────────────────────────────────────────
# PALETA
# ─────────────────────────────────────────────────────────────────────────────
ESCURO  = RGBColor(0x25, 0x30, 0x3C)
DOURADO = RGBColor(0xED, 0x9B, 0x0F)
BRANCO  = RGBColor(0xFF, 0xFF, 0xFF)
CINZA   = RGBColor(0x60, 0x6A, 0x7A)

SL_W = Inches(20)
SL_H = Inches(11.25)
MG   = Inches(0.55)

# ─────────────────────────────────────────────────────────────────────────────
# LOCALIZAR TEMPLATE
# ─────────────────────────────────────────────────────────────────────────────

def _normalizar_nome_arquivo(nome: str) -> str:
    nome = unicodedata.normalize("NFKD", str(nome))
    nome = "".join(ch for ch in nome if not unicodedata.combining(ch))
    nome = nome.lower()
    nome = re.sub(r"[^a-z0-9]+", " ", nome)
    return re.sub(r"\s+", " ", nome).strip()


def _encontrar_template() -> str | None:
    env_path = os.getenv("TEMPLATE_PPT_PATH", "").strip()
    if env_path:
        p_env = Path(env_path).expanduser()
        if p_env.exists():
            return str(p_env)

    base_dir = Path(__file__).resolve().parent
    candidatos_prioritarios = [
        base_dir / "Carnaval Sapucaí - Sesc.pptx",
        base_dir.parent / "Carnaval Sapucaí - Sesc.pptx",
        base_dir / "exemplos" / "Carnaval Sapucaí - Sesc.pptx",
        base_dir / "Carnaval_Sapucai_-_Sesc.pptx",
        base_dir.parent / "Carnaval_Sapucai_-_Sesc.pptx",
        base_dir / "exemplos" / "Carnaval_Sapucai_-_Sesc.pptx",
    ]
    for p in candidatos_prioritarios:
        if p.exists():
            return str(p)

    chaves = ("carnaval", "sapucai", "sesc")
    for area in [base_dir, base_dir.parent, base_dir / "exemplos"]:
        if not area.exists():
            continue
        for arquivo in area.rglob("*.pptx"):
            nome_norm = _normalizar_nome_arquivo(arquivo.name)
            if all(chave in nome_norm for chave in chaves):
                return str(arquivo.resolve())
    candidatos = [
        Path(__file__).parent / "Carnaval_Sapucaí_-_Sesc.pptx",
        Path(__file__).parent / "template_sesc.pptx",
        Path(__file__).parent / "template.pptx",
    ]
    for p in candidatos:
        if p.exists():
            return str(p)
    return None


# ─────────────────────────────────────────────────────────────────────────────
# CLONAGEM DE SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def _clonar_slide(prs_dest: Presentation, slide_origem) -> object:
    """
    Clona um slide de outra apresentação para prs_dest.
    Copia o XML completo e todos os relacionamentos (imagens, etc).
    """
    # Layout blank
    layout = prs_dest.slide_layouts[6]
    novo   = prs_dest.slides.add_slide(layout)

    # Limpar shapes existentes
    sp_tree = novo.shapes._spTree
    for el in list(sp_tree):
        tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            sp_tree.remove(el)

    # Copiar XML do slide origem
    origem_tree = slide_origem.shapes._spTree
    for el in origem_tree:
        tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            sp_tree.append(copy.deepcopy(el))

    # Copiar relacionamentos (imagens embutidas)
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    for rel in slide_origem.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_part = rel.target_part
                novo.part.relate_to(img_part, rel.reltype)
            except Exception:
                pass

    # Copiar background se existir
    try:
        bg_origem = slide_origem.background
        bg_novo   = novo.background
        if bg_origem.fill.type:
            bg_novo.fill.solid()
            bg_novo.fill.fore_color.rgb = bg_origem.fill.fore_color.rgb
    except Exception:
        pass

    return novo


# ─────────────────────────────────────────────────────────────────────────────
# HELPERS DE EDIÇÃO DE TEXTO
# ─────────────────────────────────────────────────────────────────────────────

def _set_text(shape, texto: str, size_pt: float = None,
              color: RGBColor = None, bold: bool = None):
    """Substitui o texto de um shape preservando a formatação base."""
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    # Limpar parágrafos mantendo o primeiro
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    para = tf.paragraphs[0]
    # Limpar runs
    for r in para.runs:
        r._r.getparent().remove(r._r)
    run = para.add_run()
    run.text = texto
    if size_pt:
        run.font.size = Pt(size_pt)
    if color:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold


def _find_shape(slide, nome: str):
    for sh in slide.shapes:
        if sh.name == nome:
            return sh
    return None


def _find_shape_parcial(slide, substring: str):
    """Encontra shape pelo nome parcial."""
    for sh in slide.shapes:
        if substring.lower() in sh.name.lower():
            return sh
    return None


def _find_chart(slide):
    for sh in slide.shapes:
        if sh.shape_type == 3:  # CHART
            return sh
    return None


# ─────────────────────────────────────────────────────────────────────────────
# ATUALIZAR GRÁFICO
# ─────────────────────────────────────────────────────────────────────────────

def _atualizar_grafico(slide, cats: list, vals: list,
                       cor_barra: RGBColor = None):
    """
    Substitui o gráfico existente no slide por um novo com os dados fornecidos.
    Mantém posição e tamanho do gráfico original.
    """
    cor_barra = cor_barra or ESCURO

    # Encontrar gráfico existente
    chart_shape = _find_chart(slide)

    if chart_shape:
        l, t, w, h = (chart_shape.left, chart_shape.top,
                      chart_shape.width, chart_shape.height)
        # Remover gráfico antigo
        sp = chart_shape._element
        sp.getparent().remove(sp)
    else:
        l = MG
        t = Inches(2.0)
        w = SL_W - 2 * MG
        h = Inches(8.2)

    # Criar novo gráfico
    if not cats:
        return

    cd = ChartData()
    cd.categories = cats
    cd.add_series("", vals)

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd)
    chart = cf.chart

    chart.has_legend = False
    chart.has_title  = False

    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = cor_barra

    series.data_labels.show_value     = True
    series.data_labels.number_format  = "0.0%"
    series.data_labels.font.size      = Pt(14)
    series.data_labels.font.bold      = True
    series.data_labels.font.color.rgb = BRANCO

    try:
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.size      = Pt(13)
        cat_ax.tick_labels.font.color.rgb = ESCURO

        val_ax = chart.value_axis
        val_ax.visible       = False
        val_ax.maximum_scale = min(max(vals) * 1.18, 1.0) if vals else 1.0

        chart.plot_area.format.fill.background()
        chart.chart_area.format.fill.background()
    except Exception:
        pass


# ─────────────────────────────────────────────────────────────────────────────
# NUVEM DE PALAVRAS (sem dependências externas)
# ─────────────────────────────────────────────────────────────────────────────

def _nuvem_palavras(slide, tabela, l, t, w, h):
    """Nuvem de palavras proporcional à frequência."""
    import pandas as pd, random

    dados = tabela[
        (~tabela[" "].str.lower().str.strip().isin(["total"])) &
        (~tabela.get("is_sub", pd.Series([False]*len(tabela))).fillna(False))
    ].nlargest(40, "Total").copy()

    if dados.empty:
        return

    max_v = dados["Total"].max()
    min_v = dados["Total"].min()
    rng   = max(max_v - min_v, 1)

    def fs(v): return 10 + ((v - min_v) / rng) * 34

    cores = [ESCURO, DOURADO, RGBColor(0x1A, 0x56, 0x9E), CINZA]
    random.seed(42)

    palavras = [(str(r[" "]), fs(r["Total"])) for _, r in dados.iterrows()]
    grandes  = [p for p in palavras if p[1] >= 30]
    medias   = [p for p in palavras if 18 <= p[1] < 30]
    pequenas = [p for p in palavras if p[1] < 18]
    random.shuffle(medias); random.shuffle(pequenas)
    todas = grandes + medias + pequenas

    n_lin   = 8
    h_lin   = h / n_lin
    por_lin = [[] for _ in range(n_lin)]
    for i, p in enumerate(todas):
        por_lin[i % n_lin].append(p)

    for li, lin in enumerate(por_lin):
        if not lin: continue
        n    = len(lin)
        step = w / (n + 1)
        for wi, (palavra, fsize) in enumerate(lin):
            px = l + step * (wi + 1) - Inches(0.6)
            py = t + h_lin * li + (h_lin - Pt(fsize).emu) / 2
            ww = Inches(max(len(palavra) * fsize * 0.013, 1.2))
            cor = cores[(li + wi) % len(cores)]
            try:
                txb = slide.shapes.add_textbox(px, py, ww, Inches(0.6))
                tf  = txb.text_frame
                tf.word_wrap = False
                p_obj = tf.paragraphs[0]
                p_obj.alignment = PP_ALIGN.CENTER
                run = p_obj.add_run()
                run.text           = palavra
                run.font.size      = Pt(fsize)
                run.font.bold      = fsize >= 28
                run.font.color.rgb = cor
                run.font.name      = "Calibri"
            except Exception:
                pass


# ─────────────────────────────────────────────────────────────────────────────
# PAINEL LATERAL DE OUTROS
# ─────────────────────────────────────────────────────────────────────────────

def _painel_outros(slide, tabela, l, t, w, h):
    import pandas as pd

    subs = tabela[tabela.get("is_sub",
                  pd.Series([False]*len(tabela))).fillna(False)]
    if subs.empty:
        return

    # Header
    _add_rect(slide, l, t, w, Inches(0.55), ESCURO)
    _add_txt(slide, "Outro — categorias",
             l + Inches(0.1), t + Inches(0.06),
             w - Inches(0.2), Inches(0.46),
             13, True, DOURADO)

    y      = t + Inches(0.60)
    h_item = Inches(0.46)
    max_i  = int((h - Inches(0.65)) / h_item)

    for i, (_, row) in enumerate(subs.iterrows()):
        if i >= max_i: break
        lbl = str(row[" "])[:42]
        pv  = row["%"]
        ps  = f"{float(pv)*100:.1f}%" if isinstance(pv, float) else str(pv)

        bg = RGBColor(0xEB, 0xEF, 0xF6) if i % 2 == 0 else BRANCO
        _add_rect(slide, l, y, w, h_item - Inches(0.03), bg)
        _add_txt(slide, lbl, l + Inches(0.1), y + Inches(0.05),
                 w * 0.72, h_item - Inches(0.08), 11, False,
                 RGBColor(0x44, 0x50, 0x60), italic=True)
        _add_txt(slide, ps, l + w * 0.73, y + Inches(0.05),
                 w * 0.25, h_item - Inches(0.08), 11, True, ESCURO,
                 align=PP_ALIGN.RIGHT, italic=True)
        y += h_item

    resto = len(subs) - max_i
    if resto > 0:
        _add_txt(slide, f"+ {resto} categorias…",
                 l, y, w, Inches(0.4), 10, False, CINZA, italic=True)


def _add_rect(slide, l, t, w, h, fill, line=None, lw=0):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def _add_txt(slide, text, l, t, w, h, size=14, bold=False,
             color=ESCURO, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = "Calibri"
    return txb


# ─────────────────────────────────────────────────────────────────────────────
# GERADOR PRINCIPAL
# ─────────────────────────────────────────────────────────────────────────────

def gerar_ppt(df, perguntas: list[dict],
              saida: str,
              titulo: str    = "Pesquisa IFec RJ",
              subtitulo: str = "",
              periodo: str   = "",
              secoes: dict   = None,
              template_path: str = None):
    """
    Gera PPT idêntico ao template Sesc/IFec.

    Usa o arquivo 'Carnaval_Sapucaí_-_Sesc.pptx' (ou template_path)
    como base — clona os slides decorativos e substitui apenas os dados.
    """
    from tabulador import tabular_pergunta
    import pandas as pd

    # Localizar template
    tmpl_path = template_path or _encontrar_template()
    if not tmpl_path:
        print("[PPT] Template não encontrado — usando fallback simples.")
        return _gerar_fallback(df, perguntas, saida, titulo, subtitulo,
                               periodo, secoes)

    prs_tmpl = Presentation(tmpl_path)
    prs_dest = Presentation()
    prs_dest.slide_width  = SL_W
    prs_dest.slide_height = SL_H

    # Índices dos slides template (0-based)
    # 0 = capa, 1 = objetivo, 2 = sumário, 3 = seção metodologia,
    # 4 = metodologia (cards), 5 = seção perfil, 6 = perfil,
    # 11 = slide resultado padrão (1 gráfico),
    # 12 = slide resultado com 2 cards (RU + RM lado a lado)
    IDX_CAPA       = 0
    IDX_OBJETIVO   = 1
    IDX_SECAO      = 3    # slide de seção com foto
    IDX_METODOLOGIA= 4    # slide com 4 cards
    IDX_RESULTADO  = 11   # slide resultado padrão

    ativas = [p for p in perguntas
              if p.get("ativo", True) and p["tipo"] != "IGNORAR"]

    # Separar perfil de resultado
    _kw_perfil = {"gênero","genero","escolaridade","estado civil","filhos",
                  "renda","moradores","trabalhando","setor","idade","perfil",
                  "quantas pessoas","quantos filhos"}
    pergs_perfil    = [p for p in ativas
                       if any(k in p["pergunta"].lower() for k in _kw_perfil)]
    pergs_resultado = [p for p in ativas if p not in pergs_perfil]

    if secoes is None:
        secoes = {}
        if pergs_resultado:
            secoes["Resultados"] = [p["num"] for p in pergs_resultado]
        if pergs_perfil:
            secoes["Perfil dos Entrevistados"] = [p["num"] for p in pergs_perfil]

    idx_p = {p["num"]: p for p in ativas}
    num   = 1

    # ── CAPA ─────────────────────────────────────────────────────────────────
    slide = _clonar_slide(prs_dest, prs_tmpl.slides[IDX_CAPA])
    # Atualizar título
    for sh in slide.shapes:
        if sh.name == 'Título 1' and hasattr(sh, 'text_frame'):
            _set_text(sh, titulo, size_pt=57.6, color=ESCURO, bold=True)
        elif 'Texto 3' in sh.name and hasattr(sh, 'text_frame'):
            if 'Realização' in sh.text_frame.text:
                _set_text(sh, subtitulo or "Realização IFec RJ",
                          size_pt=43.2, color=DOURADO)
            elif 'Fevereiro' in sh.text_frame.text or not subtitulo:
                _set_text(sh, periodo or "",
                          size_pt=28, color=DOURADO)
        elif sh.name == 'Espaço Reservado para Número de Slide 62':
            _set_text(sh, str(num))
    num += 1

    # ── SEÇÕES E RESULTADOS ──────────────────────────────────────────────────
    for i_sec, (nome_sec, nums_sec) in enumerate(secoes.items(), 1):
        # Slide de seção
        slide = _clonar_slide(prs_dest, prs_tmpl.slides[IDX_SECAO])
        for sh in slide.shapes:
            if sh.name == 'TextBox 30' and hasattr(sh, 'text_frame'):
                _set_text(sh, nome_sec, size_pt=71, color=ESCURO, bold=True)
            elif 'Número de Slide' in sh.name:
                _set_text(sh, str(num))
            elif sh.name == 'TextBox 26' and hasattr(sh, 'text_frame'):
                _set_text(sh, str(i_sec), size_pt=138, color=BRANCO, bold=True)
        num += 1

        # Slide de metodologia após seção 1
        if i_sec == 1:
            slide = _clonar_slide(prs_dest, prs_tmpl.slides[IDX_METODOLOGIA])
            for sh in slide.shapes:
                if 'Número de Slide' in sh.name:
                    _set_text(sh, str(num))
            num += 1

        # Slides de resultado
        for num_p in nums_sec:
            p = idx_p.get(num_p)
            if not p:
                continue
            try:
                tabela = tabular_pergunta(df, p)
                if tabela.empty:
                    continue
            except Exception as e:
                print(f"[PPT] Erro ao tabular {num_p}: {e}")
                continue

            # Clonar slide resultado
            slide = _clonar_slide(prs_dest, prs_tmpl.slides[IDX_RESULTADO])

            # Título
            for sh in slide.shapes:
                if sh.name == 'Título 42' and hasattr(sh, 'text_frame'):
                    titulo_txt = p["pergunta"]
                    if len(titulo_txt) > 80:
                        titulo_txt = titulo_txt[:77] + "..."
                    _set_text(sh, titulo_txt.upper(),
                              size_pt=54, color=DOURADO, bold=True)
                elif 'Número de Slide' in sh.name:
                    _set_text(sh, str(num))
                elif sh.name == 'CaixaDeTexto 52' and hasattr(sh, 'text_frame'):
                    # Nota de base
                    total_row = tabela[tabela[" "].str.lower().str.strip() == "total"]
                    base_n = int(total_row["Total"].iloc[0]) if not total_row.empty else "?"
                    tipo_txt = {"RU": "Estimulada — RU", "RM": "Estimulada — RM",
                                "ABERTA": "Espontânea — RA"}.get(p["tipo"], "")
                    nota = f"{p['num']}. {p['pergunta'][:80]}  ({tipo_txt} — Base: {base_n})"
                    _set_text(sh, nota, size_pt=14)

            # Dados para o gráfico
            tem_sub = ("is_sub" in tabela.columns and tabela["is_sub"].any())
            rotulos = tabela[" "].fillna("").astype(str)
            dados = tabela[
                (~rotulos.str.lower().str.strip().isin(["total"])) &
                (~rotulos.str.startswith("NPS =")) &
                (~tabela.get("is_sub",
                  pd.Series([False]*len(tabela))).fillna(False))
            ].copy()

            cats   = dados[" "].tolist()
            valores = []
            for v in dados["%"]:
                try:   valores.append(float(v))
                except: valores.append(0.0)

            if p["tipo"] == "ABERTA":
                # Remover gráfico, colocar nuvem
                chart_sh = _find_chart(slide)
                if chart_sh:
                    l, t, w, h = (chart_sh.left, chart_sh.top,
                                  chart_sh.width, chart_sh.height)
                    chart_sh._element.getparent().remove(chart_sh._element)
                    _nuvem_palavras(slide, tabela, l, t, w, h)
            elif cats:
                if tem_sub:
                    chart_sh = _find_chart(slide)
                    if chart_sh:
                        l, t, w, h = (chart_sh.left, chart_sh.top,
                                      chart_sh.width, chart_sh.height)
                        chart_sh._element.getparent().remove(chart_sh._element)
                        gw = int(w * 0.64)
                        _atualizar_grafico_direto(
                            slide, cats, valores, l, t, gw, h)
                        _painel_outros(slide, tabela,
                                       l + gw + Inches(0.25), t,
                                       w - gw - Inches(0.25), h)
                    else:
                        _atualizar_grafico(slide, cats, valores)
                else:
                    _atualizar_grafico(slide, cats, valores)

            num += 1

    # ── ENCERRAMENTO ─────────────────────────────────────────────────────────
    # Clonar capa e adaptar para encerramento
    slide = _clonar_slide(prs_dest, prs_tmpl.slides[IDX_CAPA])
    for sh in slide.shapes:
        if sh.name == 'Título 1' and hasattr(sh, 'text_frame'):
            _set_text(sh, "Obrigado!", size_pt=57.6, color=ESCURO, bold=True)
        elif 'Texto 3' in sh.name and hasattr(sh, 'text_frame'):
            if 'Realização' in sh.text_frame.text:
                _set_text(sh, "IFec RJ — Gerência de Processos e Projetos",
                          size_pt=28, color=DOURADO)
            else:
                _set_text(sh, "Coordenação de Processos e Inteligência de Negócios",
                          size_pt=22, color=DOURADO)
        elif 'Número de Slide' in sh.name:
            _set_text(sh, str(num))

    prs_dest.save(saida)
    print(f"PPT gerado: {saida}  ({num} slides)")
    return saida


def _atualizar_grafico_direto(slide, cats, vals, l, t, w, h,
                               cor=None):
    """Adiciona gráfico diretamente sem procurar shape existente."""
    cor = cor or ESCURO
    if not cats: return

    cd = ChartData()
    cd.categories = cats
    cd.add_series("", vals)

    cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd)
    chart = cf.chart
    chart.has_legend = False
    chart.has_title  = False

    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = cor
    series.data_labels.show_value     = True
    series.data_labels.number_format  = "0.0%"
    series.data_labels.font.size      = Pt(14)
    series.data_labels.font.bold      = True
    series.data_labels.font.color.rgb = BRANCO

    try:
        chart.category_axis.tick_labels.font.size      = Pt(13)
        chart.category_axis.tick_labels.font.color.rgb = ESCURO
        chart.value_axis.visible       = False
        chart.value_axis.maximum_scale = min(max(vals) * 1.18, 1.0) if vals else 1.0
        chart.plot_area.format.fill.background()
        chart.chart_area.format.fill.background()
    except Exception:
        pass


# ─────────────────────────────────────────────────────────────────────────────
# FALLBACK — gera PPT simples se template não encontrado
# ─────────────────────────────────────────────────────────────────────────────

def _gerar_fallback(df, perguntas, saida, titulo, subtitulo, periodo, secoes):
    """PPT simples sem template — usado se o arquivo original não existir."""
    from tabulador import tabular_pergunta
    import pandas as pd

    prs = Presentation()
    prs.slide_width  = SL_W
    prs.slide_height = SL_H

    def blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def bg(slide, c):
        f = slide.background.fill; f.solid(); f.fore_color.rgb = c

    ativas = [p for p in perguntas
              if p.get("ativo", True) and p["tipo"] != "IGNORAR"]
    num = 1

    # Capa
    slide = blank(); bg(slide, ESCURO)
    _add_txt(slide, titulo, Inches(1), Inches(2), Inches(18), Inches(4),
             50, True, BRANCO)
    _add_txt(slide, subtitulo, Inches(1), Inches(6.5), Inches(18), Inches(1),
             26, False, DOURADO)
    _add_txt(slide, str(num), Inches(17.5), Inches(10.5), Inches(2), Inches(0.5),
             16, False, CINZA, PP_ALIGN.RIGHT)
    num += 1

    for p in ativas:
        try:
            tabela = tabular_pergunta(df, p)
            if tabela.empty: continue
        except Exception:
            continue

        slide = blank(); bg(slide, RGBColor(0xF4, 0xF6, 0xF9))

        _add_txt(slide, p["pergunta"].upper(),
                 Inches(0.55), Inches(0.15), Inches(18.9), Inches(1.0),
                 28, True, DOURADO, PP_ALIGN.CENTER)

        dados = tabela[
            (~tabela[" "].str.lower().str.strip().isin(["total"])) &
            (~tabela.get("is_sub", pd.Series([False]*len(tabela))).fillna(False))
        ]
        cats  = dados[" "].tolist()[:10]
        vals  = []
        for v in dados["%"].tolist()[:10]:
            try: vals.append(float(v))
            except: vals.append(0.0)

        if cats:
            cd = ChartData()
            cd.categories = cats
            cd.add_series("", vals)
            cf = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.55), Inches(1.3),
                SL_W - Inches(1.1), SL_H - Inches(2.1), cd)
            chart = cf.chart
            chart.has_legend = False; chart.has_title = False
            s = chart.series[0]
            s.format.fill.solid(); s.format.fill.fore_color.rgb = ESCURO
            s.data_labels.show_value = True
            s.data_labels.number_format = "0.0%"
            s.data_labels.font.size = Pt(14)
            s.data_labels.font.bold = True
            s.data_labels.font.color.rgb = BRANCO
            try:
                chart.value_axis.visible = False
                chart.plot_area.format.fill.background()
                chart.chart_area.format.fill.background()
            except Exception:
                pass

        _add_txt(slide, str(num), Inches(17.5), Inches(10.5),
                 Inches(2), Inches(0.5), 16, False, CINZA, PP_ALIGN.RIGHT)
        num += 1

    prs.save(saida)
    print(f"PPT (fallback) gerado: {saida}  ({num} slides)")
    return saida
